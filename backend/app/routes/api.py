"""HTTP endpoints for the REACH dashboard.

GET  /api/health
GET  /api/filters            (country?, state?)
GET  /api/kpis               (country?, state?, lga?, year?, quarter?, round?)
GET  /api/by-country         (year?, quarter?, round?)
GET  /api/by-state           (country, state?, year?, quarter?, round?)
GET  /api/trend              (country?, state?, lga?, year?, round?)
GET  /api/boundaries         (country?, state?, admin_level=1|2, year?, quarter?, round?)
GET  /api/download           (country?, state?, lga?, year?, quarter?, round?, format=csv|pdf|pptx)
"""
from __future__ import annotations

import csv
import io
import json
import sqlite3
from typing import Optional

from fastapi import APIRouter, HTTPException, Query
from fastapi.responses import JSONResponse, StreamingResponse, Response
from shapely.geometry import shape, mapping
from shapely.ops import unary_union

from ..database import get_conn

router = APIRouter()

VALID_COUNTRIES = {"NIGERIA", "NIGER", "MALI"}


def _norm(country: Optional[str]) -> Optional[str]:
    if not country:
        return None
    c = country.strip().upper()
    if c in {"ALL", "REGIONAL", ""}:
        return None
    if c not in VALID_COUNTRIES:
        raise HTTPException(status_code=400, detail=f"Unknown country: {country}")
    return c


def _where(
    country: Optional[str],
    year: Optional[int],
    quarter: Optional[str],
    round_: Optional[int],
    state: Optional[str] = None,
    lga: Optional[str] = None,
):
    clauses, params = [], []
    if country:
        clauses.append("country = ?"); params.append(country)
    if year is not None:
        clauses.append("year = ?"); params.append(year)
    if quarter:
        clauses.append("quarter_label = ?"); params.append(quarter)
    if round_ is not None:
        clauses.append("rounds = ?"); params.append(round_)
    if state:
        clauses.append("spatial_state = ?"); params.append(state.upper())
    if lga:
        clauses.append("lga = ?"); params.append(lga.upper())
    where = (" WHERE " + " AND ".join(clauses)) if clauses else ""
    return where, params


def round_pct(num: int, den: int) -> int:
    """Coverage percentages are displayed as whole integers per UI spec."""
    if not den:
        return 0
    return int(round((num / den) * 100))


# ---------------------------------------------------------------------------
@router.get("/health")
def health():
    try:
        with get_conn(read_only=True) as conn:
            rows = conn.execute("SELECT COUNT(*) AS n FROM treatments").fetchone()["n"]
            polys = conn.execute("SELECT COUNT(*) AS n FROM boundaries").fetchone()["n"]
    except sqlite3.OperationalError:
        rows, polys = -1, -1
    return {"status": "ok", "rows": rows, "boundaries": polys}


# ---------------------------------------------------------------------------
def _quarter_sort_key(q):
    parts = (q or "").split()
    if len(parts) == 2:
        if parts[0].startswith("Q"):
            try: return (int(parts[1]), int(parts[0][1:]))
            except ValueError: return (9999, 9)
        try: return (int(parts[0]), int(parts[1][1:]))
        except ValueError: return (9999, 9)
    return (9999, 9)


@router.get("/filters")
def filters(
    country: Optional[str] = None,
    state: Optional[str] = None,
    year: Optional[int] = None,
    quarter: Optional[str] = None,
    round: Optional[int] = Query(None, alias="round"),
):
    """Cascading filter options — each list reflects values that are still
    valid given the OTHER active filters (so picking Year=2026 trims the
    Round list to rounds that exist in 2026, etc.). Quarters are returned
    in reverse-chronological order (latest first)."""
    c = _norm(country)
    f = {"country": c, "state": state, "year": year, "quarter": quarter, "round_": round}
    def _opts(exclude):
        g = dict(f); g[exclude] = None
        return _where(g["country"], g["year"], g["quarter"], g["round_"], state=g["state"])

    with get_conn(read_only=True) as conn:
        w, p = _opts("year")
        years = sorted({r[0] for r in conn.execute(
            f"SELECT DISTINCT year FROM treatments{w}", p).fetchall() if r[0] is not None})

        w, p = _opts("quarter")
        quarters = [r[0] for r in conn.execute(
            f"SELECT DISTINCT quarter_label FROM treatments{w}", p).fetchall() if r[0]]
        quarters.sort(key=_quarter_sort_key, reverse=True)

        w, p = _opts("round_")
        rounds_ = sorted({r[0] for r in conn.execute(
            f"SELECT DISTINCT rounds FROM treatments{w}", p).fetchall() if r[0] is not None})

        w, p = _opts("state")
        states = sorted({r[0] for r in conn.execute(
            f"SELECT DISTINCT spatial_state FROM treatments{w}", p).fetchall() if r[0]})

        lgas = []
        if c:
            w, p = _where(c, year, quarter, round, state=state)
            lgas = sorted({r[0] for r in conn.execute(
                f"SELECT DISTINCT lga FROM treatments{w}", p).fetchall() if r[0]})

    return {"years": years, "quarters": quarters, "rounds": rounds_, "states": states, "lgas": lgas}


# ---------------------------------------------------------------------------
@router.get("/kpis")
def kpis(
    country: Optional[str] = None,
    state: Optional[str] = None,
    lga: Optional[str] = None,
    year: Optional[int] = None,
    quarter: Optional[str] = None,
    round: Optional[int] = Query(None, alias="round"),
):
    c = _norm(country)
    where, params = _where(c, year, quarter, round, state=state, lga=lga)
    sql = f"""
        SELECT
            COALESCE(SUM(children_eligible), 0)      AS eligible,
            COALESCE(SUM(children_treated), 0)       AS treated,
            COALESCE(SUM(severe_adverse_event), 0)   AS sae,
            COALESCE(SUM(death_averted), 0)          AS death_averted
        FROM treatments{where}
    """
    with get_conn(read_only=True) as conn:
        row = conn.execute(sql, params).fetchone()
    eligible = int(row["eligible"] or 0)
    treated = int(row["treated"] or 0)
    # Deaths-averted estimate: central rate is 1 death per 1,000 treated
    # children (MORDOR-calibrated). MORDOR's 95 % CI on the mortality RR
    # (0.79–0.93) implies the *true* averted count plausibly lies in a band
    # roughly half to one-and-a-half times the central estimate.
    # Note: the function arg `round` shadows the builtin, so we avoid `round()`
    # here. Add-half-then-floor gives the same nearest-integer behaviour.
    deaths_mid  = int(treated // 1000)
    deaths_low  = int(treated * 0.5 / 1000 + 0.5)
    deaths_high = int(treated * 1.5 / 1000 + 0.5)
    return {
        "children_eligible": eligible,
        "children_treated": treated,
        "percentage_treated": round_pct(treated, eligible),
        "severe_adverse_event": int(row["sae"] or 0),
        "deaths_averted": deaths_mid,
        "deaths_averted_low":  deaths_low,
        "deaths_averted_high": deaths_high,
        "deaths_averted_note": "Central 1/1,000; range derived from MORDOR 95% CI",
    }


@router.get("/by-country")
def by_country(year: Optional[int] = None, quarter: Optional[str] = None, round: Optional[int] = Query(None, alias="round")):
    where, params = _where(None, year, quarter, round)
    sql = f"""
        SELECT country,
               COALESCE(SUM(children_eligible), 0) AS eligible,
               COALESCE(SUM(children_treated), 0)  AS treated
        FROM treatments{where}
        GROUP BY country
        ORDER BY country
    """
    with get_conn(read_only=True) as conn:
        rows = conn.execute(sql, params).fetchall()
    return [{
        "country": r["country"], "eligible": int(r["eligible"]),
        "treated": int(r["treated"]),
        "percentage": round_pct(int(r["treated"]), int(r["eligible"])),
    } for r in rows]


@router.get("/by-state")
def by_state(
    country: str = Query(...),
    state: Optional[str] = None,
    year: Optional[int] = None,
    quarter: Optional[str] = None,
    round: Optional[int] = Query(None, alias="round"),
):
    c = _norm(country)
    if c is None:
        raise HTTPException(status_code=400, detail="country is required")
    where, params = _where(c, year, quarter, round, state=state)
    group = "lga" if state else "spatial_state"
    sql = f"""
        SELECT {group} AS state,
               COALESCE(SUM(children_eligible), 0) AS eligible,
               COALESCE(SUM(children_treated), 0)  AS treated
        FROM treatments{where}
        GROUP BY {group}
        ORDER BY treated DESC
    """
    with get_conn(read_only=True) as conn:
        rows = conn.execute(sql, params).fetchall()
    return [{
        "state": r["state"], "eligible": int(r["eligible"]),
        "treated": int(r["treated"]),
        "percentage": round_pct(int(r["treated"]), int(r["eligible"])),
    } for r in rows if r["state"]]


@router.get("/trend")
def trend(
    country: Optional[str] = None,
    state: Optional[str] = None,
    lga: Optional[str] = None,
    year: Optional[int] = None,
    round: Optional[int] = Query(None, alias="round"),
):
    c = _norm(country)
    where, params = _where(c, year, None, round, state=state, lga=lga)
    sql = f"""
        SELECT quarter_label AS quarter,
               COALESCE(SUM(children_eligible), 0) AS eligible,
               COALESCE(SUM(children_treated), 0)  AS treated
        FROM treatments{where}
        GROUP BY quarter_label
    """
    with get_conn(read_only=True) as conn:
        rows = conn.execute(sql, params).fetchall()
    def sort_key(q):
        parts = (q or "").split()
        if len(parts) == 2 and parts[0].startswith("Q"):
            try: return (int(parts[1]), int(parts[0][1:]))
            except ValueError: return (9999, 9)
        return (9999, 9)
    return sorted(
        ({
            "quarter": r["quarter"],
            "eligible": int(r["eligible"]),
            "treated": int(r["treated"]),
            "percentage": round_pct(int(r["treated"]), int(r["eligible"])),
        } for r in rows if r["quarter"]),
        key=lambda p: sort_key(p["quarter"]),
    )


# ---------------------------------------------------------------------------
def _round_growth_per_state(conn, country):
    """Per-state growth = (treated_latest_round − treated_prev_round)/treated_prev."""
    rows = conn.execute("""
        SELECT spatial_state AS s, rounds AS r,
               COALESCE(SUM(children_treated),0) AS t
        FROM treatments WHERE country=? GROUP BY spatial_state, rounds
    """, (country,)).fetchall()
    by_state = {}
    for r in rows:
        if not r["s"]: continue
        by_state.setdefault(r["s"], []).append((int(r["r"]), int(r["t"])))
    out = {}
    for s, series in by_state.items():
        series.sort()
        if len(series) < 2:
            out[s.upper()] = {"growth_pct": None, "treated_prev": series[-1][1] if series else 0,
                              "treated_now": series[-1][1] if series else 0, "has_growth": False}
            continue
        prev, now = series[-2][1], series[-1][1]
        gp = ((now - prev) / prev * 100) if prev else None
        out[s.upper()] = {
            "growth_pct": round(gp, 1) if gp is not None else None,
            "treated_prev": prev, "treated_now": now,
            "round_prev": series[-2][0], "round_now": series[-1][0],
            "has_growth": gp is not None,
        }
    return out


# Lazy in-memory cache of admin-0 (country outline) GeoJSON. Computed from
# the union of the admin-1 polygons; never changes between requests.
_COUNTRY_OUTLINE_CACHE: dict[str, dict] = {}

# Cache of welded has-data shapes keyed by (country, filter-signature).
_DATA_BLOB_CACHE: dict[tuple, dict] = {}


def _data_blob_geojson(conn, ct: str, year, quarter, round_) -> Optional[dict]:
    """Return ONE welded GeoJSON polygon covering all admin-1 areas with data
    for this country under the active filter set. Used by the Landing /
    Regional low-zoom view so the coloured region appears seamless — no
    inter-state slivers."""
    key = (ct, year, quarter, round_)
    cached = _DATA_BLOB_CACHE.get(key)
    if cached is not None:
        return cached
    where, params = _where(ct, year, quarter, round_)
    rows = conn.execute(
        f"""SELECT spatial_state AS s,
                   COALESCE(SUM(children_eligible),0) AS e,
                   COALESCE(SUM(children_treated),0)  AS t
              FROM treatments{where}
              GROUP BY spatial_state""", params,
    ).fetchall()
    has_data = {str(r["s"]).upper() for r in rows
                if r["s"] and (int(r["e"]) > 0 or int(r["t"]) > 0)}
    if not has_data:
        _DATA_BLOB_CACHE[key] = None
        return None
    geoms = []
    for b in conn.execute(
        "SELECT name_upper, geojson FROM boundaries WHERE country=? AND admin_level=1",
        (ct,),
    ).fetchall():
        if b["name_upper"] not in has_data:
            continue
        try:
            g = shape(json.loads(b["geojson"]))
            if not g.is_valid:
                g = g.buffer(0)
            geoms.append(g)
        except (json.JSONDecodeError, ValueError):
            continue
    if not geoms:
        _DATA_BLOB_CACHE[key] = None
        return None
    merged = unary_union(geoms)
    try:
        merged = merged.buffer(0)
        # Close gaps between independently-simplified neighbours, then
        # simplify slightly so the payload stays small.
        merged = merged.buffer(0.03).buffer(-0.025)
        merged = merged.simplify(0.04, preserve_topology=True)
    except Exception:
        pass
    out = mapping(merged)
    _DATA_BLOB_CACHE[key] = out
    return out


def _country_outline_geojson(conn, ct: str) -> Optional[dict]:
    """Return a clean GeoJSON polygon for the whole country.

    Preferred path: the admin-0 row loaded from the official GADM admin-0
    shapefile (one clean country outline, no internal lines).
    Fallback: synthesise the outline by welding the admin-1 polygons; used
    only when an admin-0 row isn't present (e.g. a DB ingested before the
    admin-0 loader was added).
    """
    cached = _COUNTRY_OUTLINE_CACHE.get(ct)
    if cached is not None:
        return cached

    # 1. Prefer the official GADM admin-0 polygon if it's in the DB.
    row = conn.execute(
        "SELECT geojson FROM boundaries WHERE country=? AND admin_level=0 LIMIT 1",
        (ct,),
    ).fetchone()
    if row:
        try:
            geom = json.loads(row["geojson"])
            _COUNTRY_OUTLINE_CACHE[ct] = geom
            return geom
        except json.JSONDecodeError:
            pass  # fall through to synthesis
    geoms = []
    for r in conn.execute(
        "SELECT geojson FROM boundaries WHERE country=? AND admin_level=1",
        (ct,),
    ).fetchall():
        try:
            g = shape(json.loads(r["geojson"]))
            if not g.is_valid:
                g = g.buffer(0)
            geoms.append(g)
        except (json.JSONDecodeError, ValueError):
            continue
    if not geoms:
        return None
    merged = unary_union(geoms)
    try:
        # Fix any leftover topology issues, then weld near-touching vertices
        # using an outward+inward buffer (clipper-style "closing").
        merged = merged.buffer(0)
        merged = merged.buffer(0.02).buffer(-0.02)
        merged = merged.simplify(0.04, preserve_topology=True)
    except Exception:
        pass

    # Drop tiny artifact polygons / slivers — anything smaller than ~0.05
    # square degrees (about the size of a small island; well below any real
    # state). Keeps Nigeria's main land mass and any genuinely large island.
    from shapely.geometry import MultiPolygon, Polygon
    AREA_FLOOR = 0.05
    if merged.geom_type == "MultiPolygon":
        kept = [p for p in merged.geoms if p.area >= AREA_FLOOR]
        if kept:
            merged = MultiPolygon(kept) if len(kept) > 1 else kept[0]
    # Same threshold for inner holes — drop tiny donuts so the shape reads as
    # one piece at a glance.
    HOLE_FLOOR = 0.05
    def _strip_small_holes(poly: Polygon) -> Polygon:
        big = [ring for ring in poly.interiors if Polygon(ring).area >= HOLE_FLOOR]
        return Polygon(poly.exterior, big)
    if merged.geom_type == "Polygon":
        merged = _strip_small_holes(merged)
    elif merged.geom_type == "MultiPolygon":
        merged = MultiPolygon([_strip_small_holes(p) for p in merged.geoms])

    result = mapping(merged)
    _COUNTRY_OUTLINE_CACHE[ct] = result
    return result


@router.get("/boundaries")
def boundaries(
    country: Optional[str] = None,
    state: Optional[str] = None,
    admin_level: int = Query(1, ge=0, le=2),
    year: Optional[int] = None,
    quarter: Optional[str] = None,
    round: Optional[int] = Query(None, alias="round"),
    view: str = Query("coverage", pattern="^(coverage|growth)$"),
    kind: str = Query("outline", pattern="^(outline|data)$"),
):
    c = _norm(country)
    countries = [c] if c else list(VALID_COUNTRIES)
    features = []

    # ----- admin_level=0 with kind=data: welded with-data shape per country -
    # Used by the Landing / Regional low-zoom view so the coloured region
    # reads as one piece, no internal slivers.
    if admin_level == 0 and kind == "data":
        with get_conn(read_only=True) as conn:
            for ct in countries:
                where, params = _where(ct, year, quarter, round)
                row = conn.execute(
                    f"""SELECT COALESCE(SUM(children_eligible),0) AS eligible,
                               COALESCE(SUM(children_treated),0)  AS treated
                          FROM treatments{where}""", params
                ).fetchone()
                elig = int(row["eligible"] or 0)
                trt  = int(row["treated"]  or 0)
                pct  = round_pct(trt, elig)
                geom = _data_blob_geojson(conn, ct, year, quarter, round)
                if not geom:
                    continue
                features.append({
                    "type": "Feature",
                    "geometry": geom,
                    "properties": {
                        "country": ct,
                        "name": ct.title(),
                        "admin_level": 0,
                        "kind": "data",
                        "has_data": True,
                        "eligible": elig,
                        "treated":  trt,
                        "percentage": pct,
                    },
                })
        return JSONResponse({"type": "FeatureCollection", "features": features})

    # ----- admin_level=0 outline: one polygon per country, coloured by % ----
    if admin_level == 0:
        with get_conn(read_only=True) as conn:
            for ct in countries:
                where, params = _where(ct, year, quarter, round)
                row = conn.execute(
                    f"""SELECT COALESCE(SUM(children_eligible),0) AS eligible,
                               COALESCE(SUM(children_treated),0)  AS treated
                          FROM treatments{where}""", params
                ).fetchone()
                elig = int(row["eligible"] or 0)
                trt  = int(row["treated"]  or 0)
                pct  = round_pct(trt, elig)
                geom = _country_outline_geojson(conn, ct)
                if not geom:
                    continue
                features.append({
                    "type": "Feature",
                    "geometry": geom,
                    "properties": {
                        "country": ct,
                        "name": ct.title(),
                        "admin_level": 0,
                        "has_data": (elig > 0 or trt > 0),
                        "eligible": elig,
                        "treated":  trt,
                        "percentage": pct,
                    },
                })
        return JSONResponse({"type": "FeatureCollection", "features": features})

    with get_conn(read_only=True) as conn:
        for ct in countries:
            metric_col = "lga" if admin_level == 2 else "spatial_state"
            where, params = _where(ct, year, quarter, round, state=state if admin_level == 2 else None)
            metric = {}
            for r in conn.execute(
                f"""SELECT {metric_col} AS name,
                           COALESCE(SUM(children_eligible),0) AS eligible,
                           COALESCE(SUM(children_treated),0)  AS treated
                    FROM treatments{where}
                    GROUP BY {metric_col}""", params).fetchall():
                if not r["name"]:
                    continue
                metric[str(r["name"]).upper()] = {
                    "eligible": int(r["eligible"]),
                    "treated": int(r["treated"]),
                    "percentage": round_pct(int(r["treated"]), int(r["eligible"])),
                }

            growth_map = _round_growth_per_state(conn, ct) if view == "growth" else {}

            q = "SELECT name, name_upper, parent, geojson FROM boundaries WHERE country=? AND admin_level=?"
            qp = [ct, admin_level]
            if admin_level == 2 and state:
                q += " AND UPPER(parent) = ?"
                qp.append(state.upper())
            for b in conn.execute(q, qp).fetchall():
                m = metric.get(b["name_upper"], {"eligible": 0, "treated": 0, "percentage": 0})
                has_data = (m["eligible"] > 0 or m["treated"] > 0)
                try:
                    geom = json.loads(b["geojson"])
                except json.JSONDecodeError:
                    continue
                props = {
                    "country": ct,
                    "name": b["name"],
                    "parent": b["parent"],
                    "admin_level": admin_level,
                    "has_data": has_data,
                    **m,
                }
                if view == "growth":
                    g = growth_map.get(b["name_upper"], {})
                    props["growth_pct"]   = g.get("growth_pct")
                    props["treated_prev"] = g.get("treated_prev", 0)
                    props["treated_now"]  = g.get("treated_now", 0)
                    props["round_prev"]   = g.get("round_prev")
                    props["round_now"]    = g.get("round_now")
                    props["has_data"]     = bool(g.get("has_growth"))
                features.append({"type": "Feature", "geometry": geom, "properties": props})
    return JSONResponse({"type": "FeatureCollection", "features": features})


# ---------------------------------------------------------------------------
def _select_filtered(conn, country, state, lga, year, quarter, round_):
    where, params = _where(country, year, quarter, round_, state=state, lga=lga)
    cur = conn.execute(f"""
        SELECT year, country, state, lga, rounds, quarter_label,
               children_eligible, children_treated, treated_fg,
               death_averted, severe_adverse_event
        FROM treatments{where}
        ORDER BY country, state, lga, year, rounds
    """, params)
    return [d[0] for d in cur.description], cur.fetchall()


def _kpi_summary(conn, country, state, lga, year, quarter, round_):
    where, params = _where(country, year, quarter, round_, state=state, lga=lga)
    row = conn.execute(f"""
        SELECT
          COALESCE(SUM(children_eligible),0) AS eligible,
          COALESCE(SUM(children_treated),0)  AS treated,
          COALESCE(SUM(severe_adverse_event),0) AS sae,
          COALESCE(SUM(death_averted),0) AS deaths
        FROM treatments{where}
    """, params).fetchone()
    elig, trt = int(row["eligible"] or 0), int(row["treated"] or 0)
    return {
        "eligible": elig, "treated": trt,
        "percentage": round_pct(trt, elig),
        "sae": int(row["sae"] or 0),
        "deaths": int(row["deaths"] or 0),
    }


def _scope_label(country, state, lga):
    parts = []
    if lga: parts.append(f"LGA: {lga}")
    if state: parts.append(f"State: {state}")
    if country: parts.append(f"Country: {country}")
    if not parts: parts.append("All Countries")
    return ", ".join(parts)


def _make_pdf(country, state, lga, year, quarter, round_) -> bytes:
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.units import mm

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=landscape(A4),
                            leftMargin=18*mm, rightMargin=18*mm,
                            topMargin=14*mm, bottomMargin=14*mm,
                            title="REACH Dashboard Report")
    styles = getSampleStyleSheet()
    title = ParagraphStyle("title", parent=styles["Title"], fontName="Helvetica-Bold",
                           fontSize=20, textColor=colors.HexColor("#111827"))
    sub = ParagraphStyle("sub", parent=styles["Normal"], fontSize=10,
                         textColor=colors.HexColor("#6b7280"))
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], fontName="Helvetica-Bold",
                        fontSize=12, textColor=colors.HexColor("#111827"))

    story = []
    story += [Paragraph("REACH Regional Dashboard — Report", title)]
    filt_bits = []
    if year:    filt_bits.append(f"Year: {year}")
    if quarter: filt_bits.append(f"Quarter: {quarter}")
    if round_:  filt_bits.append(f"Round: {round_}")
    story += [Paragraph(f"Scope: {_scope_label(country, state, lga)} &nbsp;&nbsp; "
                        f"Filters: {' · '.join(filt_bits) if filt_bits else 'none'}", sub)]
    story += [Spacer(1, 8)]

    with get_conn(read_only=True) as conn:
        s = _kpi_summary(conn, country, state, lga, year, quarter, round_)
        story += [Paragraph("Summary", h2)]
        kpi_tbl = Table([
            ["Children Eligible", "Children Treated", "Percentage Treated", "Severe Adverse Effects", "Est. Deaths Averted"],
            [f"{s['eligible']:,}", f"{s['treated']:,}", f"{s['percentage']}%", f"{s['sae']:,}", f"{s['deaths']:,}"],
        ], colWidths=[55*mm]*5)
        kpi_tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#faf1cc")),
            ("TEXTCOLOR",  (0, 0), (-1, 0), colors.HexColor("#a87f12")),
            ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",   (0, 0), (-1, 0), 9),
            ("FONTSIZE",   (0, 1), (-1, 1), 14),
            ("FONTNAME",   (0, 1), (-1, 1), "Helvetica-Bold"),
            ("ALIGN",      (0, 0), (-1, -1), "CENTER"),
            ("VALIGN",     (0, 0), (-1, -1), "MIDDLE"),
            ("BOX",        (0, 0), (-1, -1), 0.4, colors.HexColor("#e5e7eb")),
            ("INNERGRID",  (0, 0), (-1, -1), 0.4, colors.HexColor("#e5e7eb")),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING",    (0, 0), (-1, -1), 6),
        ]))
        story += [kpi_tbl, Spacer(1, 12)]

        # Detail table — limit rows so the PDF doesn't blow up
        cols, rows = _select_filtered(conn, country, state, lga, year, quarter, round_)
        story += [Paragraph(f"Records ({len(rows):,} rows — first 200 shown)", h2)]
        header = ["Year", "Country", "State", "LGA", "Round", "Quarter", "Eligible", "Treated", "% Treated", "SAE", "Deaths"]
        body = []
        for r in rows[:200]:
            d = dict(zip(cols, r))
            elig = int(d.get("children_eligible") or 0)
            trt  = int(d.get("children_treated") or 0)
            pct = round_pct(trt, elig)
            body.append([
                d.get("year"), d.get("country"), d.get("state"), d.get("lga"),
                d.get("rounds"), d.get("quarter_label"),
                f"{elig:,}", f"{trt:,}", f"{pct}%",
                d.get("severe_adverse_event"), f"{(d.get('death_averted') or 0):.0f}",
            ])
        tbl = Table([header] + body, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#3a3f47")),
            ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
            ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",   (0, 0), (-1, 0), 8.5),
            ("FONTSIZE",   (0, 1), (-1, -1), 7.5),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f6f7f9")]),
            ("BOX",        (0, 0), (-1, -1), 0.3, colors.HexColor("#e5e7eb")),
            ("INNERGRID",  (0, 0), (-1, -1), 0.3, colors.HexColor("#e5e7eb")),
            ("ALIGN",      (4, 1), (-1, -1), "RIGHT"),
            ("LEFTPADDING", (0, 0), (-1, -1), 3),
            ("RIGHTPADDING",(0, 0), (-1, -1), 3),
        ]))
        story += [tbl]

    doc.build(story)
    return buf.getvalue()


def _make_pptx(country, state, lga, year, quarter, round_) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    with get_conn(read_only=True) as conn:
        s = _kpi_summary(conn, country, state, lga, year, quarter, round_)
        cols, rows = _select_filtered(conn, country, state, lga, year, quarter, round_)

    blank = prs.slide_layouts[6]

    # Slide 1 — title
    s1 = prs.slides.add_slide(blank)
    tx = s1.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12), Inches(1.2)).text_frame
    tx.text = "REACH Regional Dashboard"
    tx.paragraphs[0].runs[0].font.size = Pt(40)
    tx.paragraphs[0].runs[0].font.bold = True
    sub = s1.shapes.add_textbox(Inches(0.6), Inches(3.8), Inches(12), Inches(1)).text_frame
    sub.text = _scope_label(country, state, lga)
    sub.paragraphs[0].runs[0].font.size = Pt(20)
    sub.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x6b, 0x72, 0x80)

    # Slide 2 — KPIs
    s2 = prs.slides.add_slide(blank)
    th = s2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.7)).text_frame
    th.text = "Summary"
    th.paragraphs[0].runs[0].font.size = Pt(28)
    th.paragraphs[0].runs[0].font.bold = True
    tiles = [
        ("Children Eligible",   f"{s['eligible']:,}"),
        ("Children Treated",    f"{s['treated']:,}"),
        ("Percentage Treated",  f"{s['percentage']}%"),
        ("Severe Adverse Effects", f"{s['sae']:,}"),
        ("Est. Deaths Averted", f"{s['deaths']:,}"),
    ]
    tile_w, tile_h = Inches(2.4), Inches(1.7)
    gap = Inches(0.15)
    start_x = Inches(0.5)
    y = Inches(1.6)
    for i, (lbl, val) in enumerate(tiles):
        x = start_x + (tile_w + gap) * i
        shp = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, tile_w, tile_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(0xfa, 0xf1, 0xcc)
        shp.line.color.rgb = RGBColor(0xec, 0xd1, 0x82)
        tf = shp.text_frame
        tf.text = val
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(24)
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x11, 0x18, 0x27)
        p = tf.add_paragraph()
        p.text = lbl
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = RGBColor(0x6b, 0x72, 0x80)

    # Slide 3 — data table (first 25 rows)
    s3 = prs.slides.add_slide(blank)
    th = s3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6)).text_frame
    th.text = f"Records ({len(rows):,} rows — first 25 shown)"
    th.paragraphs[0].runs[0].font.bold = True
    th.paragraphs[0].runs[0].font.size = Pt(20)

    header = ["Year", "Country", "State", "LGA", "Round", "Quarter", "Eligible", "Treated", "% Treated"]
    body = []
    for r in rows[:25]:
        d = dict(zip(cols, r))
        elig = int(d.get("children_eligible") or 0)
        trt  = int(d.get("children_treated") or 0)
        body.append([
            str(d.get("year")), d.get("country"), d.get("state"), d.get("lga"),
            str(d.get("rounds")), d.get("quarter_label"),
            f"{elig:,}", f"{trt:,}", f"{round_pct(trt, elig)}%",
        ])
    nrows = len(body) + 1
    tbl_shape = s3.shapes.add_table(nrows, len(header),
                                    Inches(0.4), Inches(1.1),
                                    Inches(12.5), Inches(0.32) * nrows)
    tbl = tbl_shape.table
    for j, h in enumerate(header):
        c = tbl.cell(0, j); c.text = h
        for run in c.text_frame.paragraphs[0].runs:
            run.font.bold = True; run.font.size = Pt(10); run.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x3a, 0x3f, 0x47)
    for i, row in enumerate(body, start=1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j); c.text = str(v) if v is not None else ""
            for run in c.text_frame.paragraphs[0].runs:
                run.font.size = Pt(9)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _make_docx(country, state, lga, year, quarter, round_) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm, Inches
    from docx.enum.table import WD_ALIGN_VERTICAL

    doc = Document()
    for sec in doc.sections:
        sec.left_margin = Cm(1.6); sec.right_margin = Cm(1.6)
        sec.top_margin = Cm(1.6);  sec.bottom_margin = Cm(1.6)

    title = doc.add_heading("REACH Regional Dashboard — Report", level=0)
    for run in title.runs:
        run.font.color.rgb = RGBColor(0x11, 0x18, 0x27)

    filt_bits = []
    if year:    filt_bits.append(f"Year: {year}")
    if quarter: filt_bits.append(f"Quarter: {quarter}")
    if round_:  filt_bits.append(f"Round: {round_}")
    p = doc.add_paragraph()
    p.add_run(f"Scope: {_scope_label(country, state, lga)}").italic = True
    p.add_run("    ")
    p.add_run(f"Filters: {' · '.join(filt_bits) if filt_bits else 'none'}").italic = True

    with get_conn(read_only=True) as conn:
        s = _kpi_summary(conn, country, state, lga, year, quarter, round_)
        doc.add_heading("Summary", level=2)
        tbl = doc.add_table(rows=2, cols=5)
        tbl.style = "Light Grid Accent 1"
        headers = ["Children Eligible", "Children Treated", "Percentage Treated",
                   "Severe Adverse Effects", "Est. Deaths Averted"]
        values  = [f"{s['eligible']:,}", f"{s['treated']:,}", f"{s['percentage']}%",
                   f"{s['sae']:,}", f"{s['deaths']:,}"]
        for j, h in enumerate(headers):
            c = tbl.rows[0].cells[j]; c.text = h
            for r in c.paragraphs[0].runs: r.bold = True; r.font.size = Pt(9)
        for j, v in enumerate(values):
            c = tbl.rows[1].cells[j]; c.text = v
            for r in c.paragraphs[0].runs: r.bold = True; r.font.size = Pt(14)

        cols, rows = _select_filtered(conn, country, state, lga, year, quarter, round_)
        doc.add_paragraph()
        doc.add_heading(f"Records ({len(rows):,} rows — first 80 shown)", level=2)
        header = ["Year", "Country", "State", "LGA", "Round", "Quarter",
                  "Eligible", "Treated", "% Treated"]
        tbl = doc.add_table(rows=1, cols=len(header))
        tbl.style = "Light List Accent 1"
        for j, h in enumerate(header):
            c = tbl.rows[0].cells[j]; c.text = h
            for r in c.paragraphs[0].runs: r.bold = True; r.font.size = Pt(8)
        for r in rows[:80]:
            d = dict(zip(cols, r))
            elig = int(d.get("children_eligible") or 0)
            trt  = int(d.get("children_treated") or 0)
            row = tbl.add_row().cells
            vals = [d.get("year"), d.get("country"), d.get("state"), d.get("lga"),
                    d.get("rounds"), d.get("quarter_label"),
                    f"{elig:,}", f"{trt:,}", f"{round_pct(trt, elig)}%"]
            for j, v in enumerate(vals):
                row[j].text = "" if v is None else str(v)
                for rr in row[j].paragraphs[0].runs:
                    rr.font.size = Pt(8)

    out = io.BytesIO(); doc.save(out)
    return out.getvalue()


@router.get("/download")
def download(
    country: Optional[str] = None,
    state: Optional[str] = None,
    lga: Optional[str] = None,
    year: Optional[int] = None,
    quarter: Optional[str] = None,
    round: Optional[int] = Query(None, alias="round"),
    format: str = Query("csv", pattern="^(csv|pdf|pptx|docx)$"),
):
    c = _norm(country)
    scope = (c or "all").lower()

    if format == "csv":
        with get_conn(read_only=True) as conn:
            cols, rows = _select_filtered(conn, c, state, lga, year, quarter, round)
        buf = io.StringIO()
        w = csv.writer(buf); w.writerow(cols)
        for r in rows:
            w.writerow(list(r))
        buf.seek(0)
        return StreamingResponse(
            iter([buf.getvalue()]),
            media_type="text/csv",
            headers={"Content-Disposition": f"attachment; filename=reach_{scope}.csv"},
        )

    if format == "pdf":
        data = _make_pdf(c, state, lga, year, quarter, round)
        return Response(
            content=data, media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=reach_{scope}.pdf"},
        )

    if format == "pptx":
        data = _make_pptx(c, state, lga, year, quarter, round)
        return Response(
            content=data,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename=reach_{scope}.pptx"},
        )

    if format == "docx":
        data = _make_docx(c, state, lga, year, quarter, round)
        return Response(
            content=data,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename=reach_{scope}.docx"},
        )

    raise HTTPException(status_code=400, detail="unsupported format")
